# boletin ppt
import pandas as pd
import dataframe_image as dfi
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import sys, os
import matplotlib.ticker as mtick
import io
import base64

from matplotlib.dates import DateFormatter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from django.http import HttpResponse
from PIL import Image



# portada
def portada(presentacion, mes, year, pais, nombre):
    # poner en las dimsensiones 16x9
    presentacion.slide_height = Inches(9)
    presentacion.slide_width = Inches(16)
    # primera diap

    Layout = presentacion.slide_layouts[0]
    first_slide = presentacion.slides.add_slide(Layout)
    logo = "logo-" + pais + ".png"
    picture = first_slide.shapes.add_picture('LOGOS.png', Inches(9.6), Inches(7.75), Inches(6.25), Inches(1.15))
    log = first_slide.shapes.add_picture(logo, Inches(8.5), Inches(7.75), height=Inches(1.15))
    first_slide.shapes.title.text = "RESUMEN DEL MONITOR DE SEQUÍAS PARA " + pais.upper()
    first_slide.placeholders[1].text = str(mes) + ' ' + str(year)
    first_slide.shapes.title.width = Inches(13)
    first_slide.shapes.title.left = Inches(1.5)
    first_slide.shapes.title.top = Inches(2)
    first_slide.shapes.title.height = Inches(2)
    first_slide.placeholders[1].width = Inches(12)
    first_slide.placeholders[1].left = Inches(2)
    first_slide.placeholders[1].top = Inches(4.5)
    first_slide.placeholders[1].height = Inches(2)
    nom = str(nombre) + ".pptx"
    presentacion.save(nom)


# encabezado: a partir de la diap 2
def encabezado(slide, titulo, pais):
    picture = slide.shapes.add_picture('LOGOS.png', Inches(11.5), Inches(0.15), height=Inches(0.75))
    logo = "logo-" + pais + ".png"
    log = slide.shapes.add_picture(logo, Inches(10.5), Inches(0.15), Inches(0.75))
    title1 = slide.shapes.title
    title1.width = Inches(10)
    title1.top=Inches(0.25)
    title1.left = Inches(0.5)
    title1.height=Inches(1.5)
    title1.text = titulo
    title1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def cropPNGMap(img, border="y", tipo="SPI"):
    width, height = img.size
    top = 47  # Todos los mapas empiezan en este pixel
    bottom = height  # Todos los mapas terminan en este pixel
    img_res = img.crop((0, top, width, bottom))
    img_res = imgToBlackNWhite(img_res)
    # BUSCANDO EL PRIMER PIXEL NEGRO PARA SABER DONDE RECORTAR  #https://stackoverflow.com/questions/1109422/getting-list-of-pixel-values-from-pil
    pixels = img_res.load()

    left = -1
    for y in range(height):
        for x in range(width):
            if pixels[x, y] == 0:
                left = x - 3
                break
        if left >= 0:
            break
    # y representa la fila donde está el marco del mapa, la línea donde está la primera línea negra, ya lo encontré, falta el borde derecho.
    for x in range(width - 1, 0, -1):
        if pixels[x, y] == 0:
            right = x + 3
            break
    # TENIENDO LOS LÍMITES, CORTO
    if border == "y":
        img_res = img.crop((0, top, right, bottom))
    elif border == "n":
        img_res = img.crop((left, top, right, bottom))
    elif tipo == "SPI":
        img_res = img.crop((left, top, right, bottom))
    return img_res


def imgToBlackNWhite(img):
    thresh = 10
    fn = lambda x: 255 if x > thresh else 0
    r = img.convert('L').point(fn, mode='1')
    return r


# Devuelve una imagen en formato base64 pero previamente la recorta para quitarle el encabezado. Será usado primordialmente para listar los lapsos de índices a un mes y año determinado.
def sendBase64CropPNGMap(request, idIdx, wmode, year, month):
    if request.method == 'GET':
        mCol = "{:02d}".format(month)
        modeCol = "{:02d}".format(wmode)

        strFolder = ""
        if (idIdx == 0):
            strFolder = "mon"
            strSuffix = ""
        elif (idIdx == 1):
            strFolder = "spi"
            strSuffix = "_spi"
        elif (idIdx == 2):
            strFolder = "ssmi"
            strSuffix = "_HMS"
        elif (idIdx == 3):
            strFolder = "sndvi"
            strSuffix = "_NDV"
        elif (idIdx == 4):
            strFolder = "sti"
            strSuffix = "_tmp"
        elif (idIdx == 5):
            strFolder = "spei"
            strSuffix = "_spe"

        try:
            fpath = "/var/py/castehr/data/indices/" + strFolder + "/png/" + str(
                year) + "_" + mCol + "_" + modeCol + strSuffix + ".png"
            img = Image.open(fpath)
            img = cropPNGMap(img)
            # TRANSFORMANDO LA IMAGEN A BYTES PORQUE ASÍ LO REQUIERE BASE64ENCODE
            byteIO = io.BytesIO()
            img.save(byteIO, format='PNG')
            byteArr = byteIO.getvalue()
            return HttpResponse(base64.b64encode(byteArr).decode('utf-8'), content_type='application/octet-stream')

        # with open(fpath, "rb") as f:
        # return HttpResponse(base64.b64encode(f.read()).decode('utf-8'),content_type='application/octet-stream')
        except IOError:
            return HttpResponse(status=204)


def subtitulo(presentacion, nombre_ppt, slide, left, width, top=1.5, height=0.5, texto="ANÁLISIS", fuente=20):
    subT = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    subTTF = subT.text_frame
    subTTF.paragraphs[0].font.size = Pt(fuente)
    subTTF.paragraphs[0].font.bold = True
    subTTF.text = texto
    subTTF.paragraphs[0].font.size = Pt(fuente)
    subTTF.paragraphs[0].font.bold = True
    subTTF.paragraphs[0].alignment = PP_ALIGN.CENTER
    nom = nombre_ppt + ".pptx"
    presentacion.save(nom)


def cuadro_texto(presentacion, slide, nombre_ppt, left, width, height=6.25, top=2,
                 relleno="Cuadro para realizar el análisis de la imagen presentada"):
    texto = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    textoTF = texto.text_frame
    textoTF.word_wrap = True
    textoTF.paragraphs[0].font.size = Pt(20)
    textoTF.text = relleno
    textoTF.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
    nom = nombre_ppt + ".pptx"
    presentacion.save(nom)


def analisis_placeholder(presentacion, slide, nombre_ppt, left, top, width, height_total):
    subtitulo(presentacion, nombre_ppt, slide, left, width)
    cuadro_texto(presentacion, slide, nombre_ppt, left, width)


def calcMonths(mes):
    months = ['Enero', 'Febrero', 'Marzo', 'Abril', 'Mayo', 'Junio', 'Julio', 'Agosto', 'Septiembre', 'Octubre',
              'Noviembre', 'Diciembre']
    return months[mes - 1]


def rango_mes(mes, year, n):
    if mes - n < 0:
        sub = calcMonths(mes - n)[:3] + " " + str(year - 1) + "- " + calcMonths(mes)[:3] + " " + str(year)
    elif mes - n == 0:
        sub = calcMonths(mes - n)[:3] + " " + str(year - 1) + "- " + calcMonths(mes)[:3] + " " + str(year)
    else:
        sub = calcMonths(mes - n)[:3] + " " + str(year) + "- " + calcMonths(mes)[:3] + " " + str(year)
    return sub


def colocar_imagen(presentacion, image_path, slide, nombre_ppt, xinicial=0.5, yinicial=1.5, altura=7, tipo=""):
    # sacar imagen de la carpeta, depende del path
    script_dir = sys.path[0]
    img_path = os.path.join(script_dir, image_path)

    im = Image.open(img_path)
    width, height = im.size
    top = 47  # Todos los mapas empiezan en este pixel
    bottom = height  # Todos los mapas terminan en este pixel
    if tipo == "SPI":
        img_res = im.crop((134, top, width, bottom))
        img_res.save("a.png")
    elif tipo == "":
        img_res = im.crop((0, top, width, bottom))
        img_res.save("a.png")
    else:
        img_res = im.crop((75, top, width, bottom))
        img_res.save("a.png")

    # colocar mapa en la diapositiva
    mapa = slide.shapes.add_picture("a.png", Inches(xinicial), Inches(yinicial), height=Inches(altura))
    nom = nombre_ppt + ".pptx"
    # dimensiones de la imagen en el ppt
    width = mapa.width / 914400
    height = mapa.height / 914400

    presentacion.save(nom)
    os.remove("a.png")
    return width, height, mapa


# colocar imagen y cuadro de texto en una diapositiva
def mapa_texto(pais, presentacion, slide, titulo, path, nombre_ppt,
               text="texto de prueba de relleno, lugar donde se realiza el analisis"):
    encabezado(slide, titulo, pais)
    name = nombre_ppt + ".pptx"

    w, h, m = colocar_imagen(presentacion, path, slide, nombre_ppt, xinicial=1, yinicial=1.9, altura=6.75)
    x_inicial = 1 + w
    y_inicial = 2
    base_texto = 14.5 - w
    altura = 6.5
    cuadro_texto(presentacion, slide, nombre_ppt, left=x_inicial, top=y_inicial, width=base_texto, height=altura,
                 relleno=text)

    # subtitulo

    subtitulo(presentacion, nombre_ppt, slide, x_inicial, top=1.5, width=base_texto, height=0.5)
    presentacion.save(name)


def graph_serieTiempo(direccion, nombre):
    script_dir = sys.path[0]
    csv_path = os.path.join(script_dir, direccion)
    # poner datos en un dataframe
    df = pd.read_csv(csv_path, parse_dates={"fecha": ["YYYY", "MM"]})
    datos_serie = df.fillna(0)  # llenar los vacios con ceros

    # graficas
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlabel("Fecha")
    plt.grid(linestyle='--')  # grilla para la grafica
    plt.ylim(0, 100)  # poner en y del 0-100
    SMALL_SIZE = 8
    MEDIUM_SIZE = 16
    BIGGER_SIZE = 22



    # graficar las lineas apiladas
    colores = ["#FFFF40", "#FFDE6F", "#ECAA1C", "#E72C07", "#900B02"]
    x = datos_serie["fecha"]
    y = datos_serie["D0"], datos_serie["D1"], datos_serie["D2"], datos_serie["D3"], datos_serie["D4"]
    plt.stackplot(x, y, baseline="zero", colors=colores)
    ax.yaxis.set_major_formatter(mtick.PercentFormatter())
    # labels de la leyenda
    D0_patch = mpatches.Patch(color="#FFFF40", label='D0')
    D1_patch = mpatches.Patch(color="#FFDE6F", label='D1')
    D2_patch = mpatches.Patch(color="#ECAA1C", label='D2')
    D3_patch = mpatches.Patch(color="#E72C07", label='D3')
    D4_patch = mpatches.Patch(color="#900B02", label='D4')
    plt.xticks(datos_serie["fecha"], rotation=45, ha="right")  # formato de las fechas en x


    ax.xaxis.set_major_formatter(DateFormatter("%Y-%m"))
    # colocar la caja abajo para la leyenda
    #plt.legend(handles=[D0_patch, D1_patch, D2_patch, D3_patch, D4_patch], loc='upper center',
              # bbox_to_anchor=(0.5, -0.15),
               #fancybox=True, ncol=5, fontsize=14)
    for item in ([ax.title, ax.xaxis.label, ax.yaxis.label] +
                 ax.get_xticklabels() + ax.get_yticklabels()):
        item.set_fontsize(16)
    plt.tight_layout()
    plt.savefig(nombre)


def mapa_sin_leyenda(presentacion, nombre_ppt, slide, image_path, xinicial, yinicial, alturaimagen, nom_mapa):
    script_dir = sys.path[0]
    img_path = os.path.join(script_dir, image_path)

    im = Image.open(img_path)
    img_res = cropPNGMap(im)

    img_res.save(nom_mapa)

    mapa3meses = slide.shapes.add_picture(nom_mapa, Inches(xinicial), Inches(yinicial), height=Inches(alturaimagen))
    ancho_primermapa = mapa3meses.width / 914400

    os.remove(nom_mapa)
    presentacion.save(nombre_ppt + ".pptx")
    return ancho_primermapa


def mapa_subtitulo(presentacion, nombre_ppt, slide, image_path, xinicial, yinicial, alturaimagen, texto, type="",
                   altura_sub=1.8, fuente=20):
    # colocar_imagen(image_path,presentacion,slide,nombre_ppt, xinicial=0.5,yinicial=1.5,altura=7)

    map2width, map2height, mapa6meses = colocar_imagen(presentacion, image_path, slide, nombre_ppt, xinicial,
                                                       yinicial, altura=alturaimagen, tipo=type)
    subtitulo(presentacion, nombre_ppt, slide, xinicial, map2width, top=altura_sub, texto=texto, fuente=fuente)

    return map2width


# diagrama de barras
def diagrama_barras(year, mes, path_archivo):
    script_dir = sys.path[0]
    csv_path = os.path.join(script_dir, path_archivo)
    df = pd.read_csv(csv_path, header=0, index_col=0)
    df.columns = ['Región', 'Sin Afectación', 'D0', 'D1', "D2", "D3", "D4", "D0 a D4", "D1 a D4", "D2 a D4",
                         "D3 a D4"]
    regiones = ["Arica y Parinacota","Tarapaca", "Antofagasta","Atacama", "Coquimbo",
                "Valparaiso", "Metropolitana de Santiago", "Libertador Bernardo O'Higgins", "Maule",
                "Ñuble","Bio-Bio","La Araucania", "Los Rios",
                "Los Lagos", "Aysen del Gral.Ibañez del Campo","Magallanes y Antartica Chilena"]

    #df_barras["Región"] = (df_barras["Región"].str.slice(start=9))

    df.reset_index()
    for i in range(len(regiones)):
        r = df.iat[i,0].split(" ")
        if r[1]=="Metropolitana":
            t = " ".join(r[1:])
        else:
            t = " ".join(r[2:])
        df.iloc[i, 0] = t

    df["Región"] = pd.Categorical(df["Región"], regiones)
    df_barras=df.sort_values("Región")


    # graficas
    fig, ax = plt.subplots(figsize=(10, 5))

    plt.ylim(0, 100)  # poner en y del 0-100
    # plt.xticks(df_barras["Región"], rotation=45, ha="right")  # formato de los nombres en x

    # graficar las lineas
    plt.bar(df_barras["Región"], df_barras["D4"], color="#900B02", edgecolor="black")
    plt.bar(df_barras["Región"], df_barras["D3"], color="#E72C07", bottom=df_barras["D4"], edgecolor="black")
    plt.bar(df_barras["Región"], df_barras["D2"], color="#ECAA1C", bottom=df_barras["D3 a D4"], edgecolor="black")
    plt.bar(df_barras["Región"], df_barras["D1"], color="#FFDE6F", bottom=df_barras["D2 a D4"], edgecolor="black")
    plt.bar(df_barras["Región"], df_barras["D0"], color="#FFFF40", edgecolor="black", bottom=df_barras["D1 a D4"])

    plt.bar(df_barras["Región"], df_barras["Sin Afectación"], color="#F0F0F0", edgecolor="black",
            bottom=df_barras["D0 a D4"])
    ax.set_xticks("Región")
    ax.yaxis.set_major_formatter(mtick.PercentFormatter())
    plt.xticks(df_barras["Región"], rotation=45, ha="right")  # formato de los nombres en x

    nombre = str(year) + "_" + str(mes) + "_barras.png"

    plt.tight_layout()
    plt.savefig(nombre)

    return nombre


# poner la leyenda de los colores
def leyendas_colores(slide, categoria, colores, x, y, cuadro=0.6, espaciamientox=1, espaciamientoy=1.25, x2=0,
                     definiciones=[], abreviaturas=[], fuente=20):
    yinicial = y
    xinicial = x
    for i in range(len(categoria)):
        # cuadro de color
        shapes = slide.shapes
        width = height = Inches(cuadro)
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xinicial), Inches(yinicial), width, height)
        # shadow = shape.shadow
        # shadow.inherit = False
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(colores[i][0], colores[i][1], colores[i][2])
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)

        # abreviaturas
        if len(abreviaturas) > 0:
            abrev = shape.text_frame
            a = abrev.paragraphs[0]
            r = a.add_run()
            r.text = abreviaturas[i]
            r.font.size = Pt(10)
            r.font.bold = True
            if i < 4:
                r.font.color.rgb = RGBColor(0, 0, 0)
            else:
                r.font.color.rgb = RGBColor(255, 255, 255)

        # nombre de categoria
        inicio_cat = xinicial + espaciamientox + cuadro
        textoC = slide.shapes.add_textbox(Inches(inicio_cat), Inches(yinicial), Inches(3.5), height=height)
        textoTF = textoC.text_frame
        t = textoTF.paragraphs[0]
        run = t.add_run()
        run.text = categoria[i]
        font = run.font
        font.size = Pt(fuente)
        font.bold = True
        textoTF.paragraphs[0].alignment = PP_ALIGN.LEFT
        # definicion
        if len(definiciones) > 0:
            inicio_def = xinicial + espaciamientox + cuadro + 3.7
            textoD = slide.shapes.add_textbox(Inches(inicio_def), Inches(yinicial), Inches(15 - inicio_def),
                                              height=Inches(1.5))
            textoTF = textoD.text_frame
            textoTF.word_wrap = True
            t = textoTF.paragraphs[0]
            run = t.add_run()
            run.text = definiciones[i]
            font = run.font
            font.size = Pt(fuente)
            textoTF.paragraphs[0].alignment = PP_ALIGN.LEFT
            yinicial += espaciamientoy + cuadro + 0.5
        else:
            yinicial += espaciamientoy + cuadro

        if yinicial >= 8 and len(definiciones) == 0:
            yinicial = y
            xinicial += x2


## FUNCIONES PARA CADA DIAPOSITIVA
def diap1(presentacion, mes, pais, nombre_ppt, year):
    mes_letras = calcMonths(mes)
    boletin = portada(presentacion, mes_letras, year, pais, nombre_ppt)


# CLASIFICACION DE INTENSIDAD DE SEQUIA
def diap2(presentacion, nombre_ppt, pais):
    second_layout = presentacion.slide_layouts[5]
    second_slide = presentacion.slides.add_slide(second_layout)

    encabezado(second_slide, "Clasificación de la intensidad de la Sequía", pais)

    categoria = ["Anormalmente Seco (D0)", "Sequía Moderada (D1)", "Sequía Severa (D2)", "Sequía Extrema (D3)",
                 "Sequía Excepcional (D4)"]
    definicion = [
        "Se trata de una condición de sequedad, no es una categoría de sequía. Se presenta al inicio o al final de un período de sequía.",
        "Se presentan algunos daños en los cultivos y pastos; existe un alto riesgo de incendios, bajos niveles en ríos, arroyos, embalses, abrevaderos y pozos, se sugiere restricción voluntaria en el uso del agua.",
        "Probables pérdidas en cultivos o pastos, alto riesgo de incendios es común la escasez de agua, se deben imponer restricciones en su uso.",
        "Pérdidas mayores en cultivos y pastos, el riesgo de incendios forestales es extremo, se generalizan las restricciones en el uso del agua debido a su escasez.",
        "Pérdidas excepcionales y generalizadas de cultivos o pastos, riesgo excepcional de incendios, escasez total de agua en embalses, arroyos y pozos, es probable una situación de emergencia debido a la ausencia de agua."]

    colores = [[255, 255, 0], [255, 211, 127], [230, 152, 0], [230, 0, 0], [115, 0, 0]]
    yinicial = 2
    xinicial = 1
    leyendas_colores(second_slide, categoria, colores, xinicial, yinicial,
                     definiciones=definicion, espaciamientoy=0.25)

    presentacion.save(nombre_ppt + ".pptx")


# DIAPOSITIVA 2: CONDICIONES DE SEQUIA MES ANTERIOR
def diap3(presentacion, year, mes, dir, nombre_ppt, pais):
    nombre_img = str(year) + "_0" + str(mes - 1) + "_" + "01.png"
    third_Layout = presentacion.slide_layouts[5]
    path = dir + str("\\") + nombre_img
    third_slide = presentacion.slides.add_slide(third_Layout)
    titulo = "Condiciones de sequía -" + str(calcMonths(mes - 1)) + " " + str(year)
    mapa_texto(pais, presentacion, third_slide, titulo, path, nombre_ppt)
    presentacion.save(nombre_ppt + ".pptx")


# CONDICIONES DE SEQUIA: MES
def diap4(presentacion, year, mes, dir, nombre_ppt, pais):
    third_Layout = presentacion.slide_layouts[5]
    third_slide = presentacion.slides.add_slide(third_Layout)
    nombre_img = str(year) + "_0" + str(mes) + "_" + "01.png"
    titulo = "Condiciones de sequía -" + str(calcMonths(mes)) + " " + str(year)
    direccion = dir + str("\\") + nombre_img
    subtitulo(presentacion, nombre_ppt, third_slide, 0.5, 9, 1.25, 0.5, texto="Monitoreo preliminar")
    mapa_texto(pais, presentacion, third_slide, titulo, direccion, nombre_ppt)
    presentacion.save(nombre_ppt + ".pptx")


# CONDICIONES DE SEQUIA: SERIE DE TIEMPO+TABLA
def diap5(presentacion, dir_csv, nombre_ppt, pais):
    fourth_Layout = presentacion.slide_layouts[5]
    fourth_slide = presentacion.slides.add_slide(fourth_Layout)
    encabezado(fourth_slide, "Condiciones de sequía", pais)
    subtitulo(presentacion, nombre_ppt, fourth_slide, 0.5, 9, 1.25, 0.5, texto="Serie de tiempo mensual por categoría")

    # serie de tiempo
    script_dir = sys.path[0]
    csv_path = os.path.join(script_dir, dir_csv)

    # poner datos en un dataframe
    df = pd.read_csv(csv_path, parse_dates={"Fecha": ["YYYY", "MM"]})
    df.rename(columns={'NA': 'Sin afectación'}, inplace=True)
    datos_tabla = df.fillna(0)  # llenar los vacios con ceros
    datos_tabla.insert(2, "D0 a D4", (
                datos_tabla["D0"] + datos_tabla["D1"] + datos_tabla["D2"] + datos_tabla["D3"] + datos_tabla[
            "D4"]).round(2))
    datos_tabla.insert(3, "D1 a D4", datos_tabla["D1"] + datos_tabla["D2"] + datos_tabla["D3"] + datos_tabla["D4"])
    datos_tabla.insert(4, "D2 a D4", datos_tabla["D2"] + datos_tabla["D3"] + datos_tabla["D4"])
    datos_tabla.insert(5, "D3 a D4", +datos_tabla["D3"] + datos_tabla["D4"])
    datos_tabla.insert(6, "D4 ", datos_tabla["D4"])
    #   D0"#FFFF40"    D1 "#FFDE6F"    D2 "#ECAA1C"    D3 "#E72C07"    D4"#900B02"
    tabla = datos_tabla
    tabla["Fecha"] = df['Fecha'].astype('str')

    header=([["Categorías acumuladas","Categorías acumuladas", "Categorías acumuladas" ,"Categorías acumuladas","Categorías acumuladas","Categorías acumuladas","Categorías acumuladas","Categorías segregadas","Categorías segregadas","Categorías segregadas","Categorías segregadas","Categorías segregadas"],tabla.columns.values.tolist()])
    tabla.columns=header

    r, c = tabla.shape
    for i in range(r):
        fecha = tabla["Categorías acumuladas","Fecha"][i].split("-")
        a = fecha[0]
        m = fecha[1]
        t = calcMonths(int(m)) + " " + a
        tabla.iloc[i, 0] = t

    headers = [dict(selector='th:nth-child(5n-2)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', '#FFFF40')]),
               dict(selector='th:nth-child(5n-1)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', '#FFDE6F')]),
               dict(selector='th:nth-child(5n)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', '#ECAA1C')]),
               dict(selector='th:nth-child(5n+1)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', '#E72C07')]),
               dict(selector='th:nth-child(5n+2)',
                    props=[('text-align', 'center'), ('color', 'white'), ('background-color', '#900B02')]),
               dict(selector='th:nth-child(1)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', 'white')]),
               dict(selector='th:nth-child(2)',
                    props=[('text-align', 'center'), ('color', 'black'), ('background-color', 'white')]),
               dict(selector="th, td",props=[("border-bottom","1px solid #ddd")]),
               dict(selector="td, th", props=[("border-right", "1px solid black")]),
               dict(selector="th:not(:last-child)", props=[('text-align', 'right'), ("border-bottom","1px solid black")]),
               dict(selector="th", props=[ ("border-bottom","1px solid black")])
               ]

    df = tabla.tail(3).style.set_table_styles(headers)
    df.format(decimal=".", precision=1, na_rep="0")
    df.set_properties(**{'text-align': 'center'}).hide(axis='index')
    pd.set_option('colheader_justify', 'center')
    dfi.export(df, 'tabla1.png')
    ancho = 8

    # agregar grafica
    nom = "grafica_boletin.png"
    graph_serieTiempo(dir_csv, nom)
    yinicial = 2
    serie = fourth_slide.shapes.add_picture(nom, Inches(0.25), Inches(yinicial), width=Inches(ancho))

    os.remove(nom)
    # agregar leyenda
    categoria = ["Anormalmente Seco", "Sequía Moderada", "Sequía Severa", "Sequía Extrema",
                 "Sequía Excepcional"]
    colores = [[255, 255, 0], [255, 211, 127], [230, 152, 0], [230, 0, 0], [115, 0, 0]]
    abr = ["D0", "D1", "D2", "D3", "D4"]
    subtitulo(presentacion, nombre_ppt, fourth_slide, ancho + 0.2, 2, yinicial -0.1, 0.25,
              texto="INTENSIDAD DE LA SEQUÍA", fuente=14)
    leyendas_colores(fourth_slide, categoria, colores, ancho + 0.1, yinicial + 0.35, abreviaturas=abr, cuadro=0.4,
                     espaciamientox=0.051, espaciamientoy=0.051, fuente=12)
    altura_serie = serie.height / 914400
    # tabla
    porcentajes = fourth_slide.shapes.add_picture("tabla1.png", Inches(0.5), Inches(altura_serie + yinicial + 0.2),
                                                  width=Inches(ancho))

    analisis_placeholder(presentacion, fourth_slide, nombre_ppt, left=ancho + 2.3, top=1.25, width=15.5 - ancho - 2,
                         height_total=7.25)

    os.remove("tabla1.png")
    presentacion.save(nombre_ppt + ".pptx")


# REGIONES: GRAFICO DE BARRAS
def diap6(presentacion, dir, year, mes, nombre_ppt, pais):
    fifth_Layout = presentacion.slide_layouts[5]
    fifth_slide = presentacion.slides.add_slide(fifth_Layout)
    encabezado(fifth_slide, "Porcentaje de sequía por regiones", pais)
    path = dir + str("\\")
    dir_stats = str(year) + "_0" + str(mes) + "_01_stats_L1.csv"
    yinicial = 2
    barras = diagrama_barras(year, mes, path + dir_stats)
    diagrama = fifth_slide.shapes.add_picture(barras, Inches(2), Inches(yinicial), width=Inches(10))
    os.remove(barras)
    alt_barras = diagrama.height / 914400
    categoria = ["Sin Afectación", "Anormalmente Seco", "Sequía Moderada", "Sequía Severa", "Sequía Extrema",
                 "Sequía Excepcional"]
    colores = [[240, 240, 240], [255, 255, 0], [255, 211, 127], [230, 152, 0], [230, 0, 0], [115, 0, 0]]
    abr = [" ", "D0", "D1", "D2", "D3", "D4"]

    ancho_barras = diagrama.width / 914400

    subtitulo(presentacion, nombre_ppt, fifth_slide, ancho_barras +2.1, 2, yinicial-0.1 , 0.25,
              texto="INTENSIDAD DE LA SEQUÍA", fuente=12)

    leyendas_colores(fifth_slide, categoria, colores, ancho_barras+2, yinicial + 0.35, abreviaturas=abr, cuadro=0.4,
                     espaciamientox=0.051, espaciamientoy=0.051, fuente=12)

    cuadro_texto(presentacion, fifth_slide, nombre_ppt, 1,ancho_barras, top=7.5, height=1.2)

    presentacion.save(nombre_ppt + ".pptx")


# SERIE DE MAPA: PERSISTENCIA DE SEQUIAS
def diap7(presentacion, mes, year, dir, nombre_ppt, pais):
    seventh_Layout = presentacion.slide_layouts[5]
    seventh_slide = presentacion.slides.add_slide(seventh_Layout)
    encabezado(seventh_slide, "Persistencia de Sequías", pais)
    cuadro_texto(presentacion, seventh_slide, nombre_ppt, 1.75, 8, height=1, top=8)

    mapasSPI = os.listdir(dir)
    xinicial = 1.25
    altura_mapa = 5.5
    yinicial = 2.25
    for mapa in mapasSPI:
        path = dir + str("\\") + mapa
        if mapasSPI.index(mapa) == 1:
            subt = calcMonths(mes)[:3] + " " + str(year)

            ancho = mapa_sin_leyenda(presentacion, nombre_ppt, seventh_slide, path, xinicial, yinicial, altura_mapa,
                                     "SPI_1.png")
            subtitulo(presentacion, nombre_ppt, seventh_slide, xinicial, ancho, top=yinicial - 0.5, height=0.25,
                      texto=subt, fuente=14)
            subtitulo(presentacion, nombre_ppt, seventh_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto="1 mes", fuente=14)
            xinicial += ancho + 0.15
        elif 0 < mapasSPI.index(mapa) < len(mapasSPI) - 1:

            script_dir = sys.path[0]
            img_path = os.path.join(script_dir, path)

            im = Image.open(img_path)
            img_res = cropPNGMap(im, "n")

            img_res.save("mapa.png")
            spi = seventh_slide.shapes.add_picture("mapa.png", Inches(xinicial), Inches(yinicial),
                                                   height=Inches(altura_mapa))
            ancho = spi.width / 914400
            n = int((mapa.split("_")[2]).split(".")[0])
            if n < 12:
                subt = rango_mes(mes, year, n)
            elif n == 12:
                subt = calcMonths(mes)[:3] + " " + str(year - 1) + "- " + calcMonths(mes)[:3] + " " + str(year)
            elif n > 12:
                subt = calcMonths(mes)[:3] + " " + str(year - n // 12) + "- " + calcMonths(mes)[:3] + " " + str(year)
            subtitulo(presentacion, nombre_ppt, seventh_slide, xinicial, ancho, top=yinicial - 0.5, height=0.25,
                      texto=subt, fuente=14)
            subtitulo(presentacion, nombre_ppt, seventh_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto=str(n) + " meses", fuente=14)
            xinicial += ancho + 0.15
        elif mapasSPI.index(mapa) == len(mapasSPI) - 1:
            n = int((mapa.split("_")[2]).split(".")[0])
            subt = calcMonths(mes)[:3] + " " + str(year - n // 12) + "- " + calcMonths(mes)[:3] + " " + str(year)
            ancho = mapa_subtitulo(presentacion, nombre_ppt, seventh_slide, path, xinicial, yinicial, altura_mapa,type="p", texto=subt, altura_sub=yinicial - 0.5, fuente=14)
            subtitulo(presentacion, nombre_ppt, seventh_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto=str(n) + " meses", fuente=14)
            xinicial += ancho + 0.15

        presentacion.save(nombre_ppt + ".pptx")
    os.remove("mapa.png")
    presentacion.save(nombre_ppt + ".pptx")



# SPI: MAPAS
def diap9(presentacion, nombre_ppt, dir, mes, year, pais):
    eigth_Layout = presentacion.slide_layouts[5]
    eigth_slide = presentacion.slides.add_slide(eigth_Layout)
    encabezado(eigth_slide, "Persistencia del índice estandarizado de precipitación (SPI)", pais)
    cuadro_texto(presentacion, eigth_slide, nombre_ppt, 1.75, 8, height=1, top=8)

    mapasSPI = os.listdir(dir)
    xinicial = 1
    altura_mapa = 5.5
    yinicial = 2.25
    for mapa in mapasSPI:
        path = dir + str("\\") + mapa
        if mapasSPI.index(mapa) == 0:
            subt = calcMonths(mes)[:3] + " " + str(year)

            ancho = mapa_sin_leyenda(presentacion, nombre_ppt, eigth_slide, path, xinicial, yinicial, altura_mapa,
                                     "SPI_1.png")
            subtitulo(presentacion, nombre_ppt, eigth_slide, xinicial, ancho, top=yinicial - 0.5, height=0.25,
                      texto=subt, fuente=14)
            subtitulo(presentacion, nombre_ppt, eigth_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto="1 mes", fuente=14)
            xinicial += ancho + 0.15
        elif mapasSPI.index(mapa) < len(mapasSPI) - 1:

            script_dir = sys.path[0]
            img_path = os.path.join(script_dir, path)

            im = Image.open(img_path)
            img_res = cropPNGMap(im, "n")

            img_res.save("mapa.png")
            spi = eigth_slide.shapes.add_picture("mapa.png", Inches(xinicial), Inches(yinicial),
                                                 height=Inches(altura_mapa))
            ancho = spi.width / 914400
            n = int(mapa.split("_")[2])
            if n < 12:
                subt = rango_mes(mes, year, n)
            elif n == 12:
                subt = calcMonths(mes)[:3] + " " + str(year - 1) + "- " + calcMonths(mes)[:3] + " " + str(year)
            elif n > 12:
                subt = calcMonths(mes)[:3] + " " + str(year - n // 12) + "- " + calcMonths(mes)[:3] + " " + str(year)
            subtitulo(presentacion, nombre_ppt, eigth_slide, xinicial, ancho, top=yinicial - 0.5, height=0.25,
                      texto=subt, fuente=14)
            subtitulo(presentacion, nombre_ppt, eigth_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto=str(n) + " meses", fuente=14)
            xinicial += ancho + 0.15
        elif mapasSPI.index(mapa) == len(mapasSPI) - 1:
            n = int(mapa.split("_")[2])
            subt = calcMonths(mes)[:3] + " " + str(year - n // 12) + "- " + calcMonths(mes)[:3] + " " + str(year)
            ancho = mapa_subtitulo(presentacion, nombre_ppt, eigth_slide, path, xinicial, yinicial, altura_mapa,
                                   type="SPI", texto=subt, altura_sub=yinicial - 0.5, fuente=14)
            subtitulo(presentacion, nombre_ppt, eigth_slide, xinicial, ancho, top=yinicial - 0.3, height=0.25,
                      texto=str(n) + " meses", fuente=14)
            xinicial += ancho + 0.15

        presentacion.save(nombre_ppt + ".pptx")
    os.remove("mapa.png")
    presentacion.save(nombre_ppt + ".pptx")


# LOGO DE LA INSTITUCION
def diap10(presentacion, nombre_ppt, pais):
    tenth_layout = presentacion.slide_layouts[6]
    tenth_slide = presentacion.slides.add_slide(tenth_layout)

    info = ["Monitor de sequías de Chile" + "\n",
            "Este servicio de información busca proveer a los tomadores de decisiones, planificadores, medios de comunicación, científicos y población en general, una síntesis útil y oportuna del monitoreo de las condiciones de sequía en el territorio chileno." + "\n",
            "Para más información visite: ", "meteochile.gob.cl"]
    textoI = tenth_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(13.25), Inches(2.5))
    textoTF = textoI.text_frame
    t = textoTF.paragraphs[0]
    textoTF.word_wrap = True
    for i in range(len(info) - 1):
        run = t.add_run()
        run.text = info[i]
        font = run.font
        font.size = Pt(25)
        if i == 0:
            font.bold = True
        textoTF.paragraphs[0].alignment = PP_ALIGN.LEFT
    r = t.add_run()
    r.text = info[-1]
    hlink = r.hyperlink
    hlink.address = 'http://www.meteochile.gob.cl/PortalDMC-web/index.xhtml'
    font = r.font
    font.size = Pt(25)
    tenth_slide.shapes.add_picture("LOGOS.png", Inches(4.25), Inches(5), height=Inches(1.75))
    logo = "logo-" + pais.lower() + ".png"
    tenth_slide.shapes.add_picture(logo, Inches(2.25), Inches(5), height=Inches(1.75))
    presentacion.save(nombre_ppt + ".pptx")


##ELABORACION DEL BOLETIN: debe tener una presentacion llamada ppt

"PPT"
"DATOS"
month = 6
year = 2021
pais = "chile"


def boletin_sequias_chile(year, month):
    pais = "Chile"
    direccion_mapas = r'C:\Users\NATHI\Documents\ciifen\Boletin Sequias\CHILE MONITOR JUNIO2022'
    direccion_datos = r"C:\Users\NATHI\Documents\ciifen\Boletin Sequias\national.csv"
    direccion_SPI = r"C:\Users\NATHI\Documents\ciifen\insumos\insumos\CHILE SPI JUNIO2022"
    direccion_stats = r"C:\Users\NATHI\Documents\ciifen\Boletin Sequias\stats\stats"

    name_ppt = "Boletin_Sequias_" + pais
    ppt = Presentation()
    # poner en las dimsensiones 16x9
    ppt.slide_height = Inches(9)
    ppt.slide_width = Inches(16)

    diap1(ppt, month, pais, name_ppt, year)  # portada

    diap2(ppt, name_ppt, pais)  # Clasificacion de la intensidad de la sequia

    diap3(ppt, year, month, direccion_mapas, name_ppt, pais)  # mapa mes pasado

    diap6(ppt, direccion_stats, year, month, name_ppt, pais)  # mapa regiones: condiciones de sequia

    diap5(ppt, direccion_datos, name_ppt, pais)  # mapa serie de tiempo: condiciones de sequia

    diap7(ppt, month, year, direccion_mapas, name_ppt, pais)  # serie de mapas: persistencia de sequia

    diap9(ppt, name_ppt, direccion_SPI, month, year, pais)  # mapas SPI

    diap10(ppt, name_ppt, pais)  # logo de entidad meteorologica del pais


boletin_sequias_chile(year, month)
